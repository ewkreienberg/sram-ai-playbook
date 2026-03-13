#!/usr/bin/env python3
"""
Build the SRAM AI Adoption Playbook presentation.

McKinsey-grade executive deck for an 8-minute class presentation.
Uses Duarte's "What Is / What Could Be" contrast pattern with
action titles (each slide title states the conclusion, not the topic).

Color scheme: SRAM brand palette, light theme.

Usage:
    python3 tools/build_presentation.py

Output:
    deliverables/SRAM_AI_Adoption_Playbook.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- SRAM brand palette, light theme ---
# Sourced from sram.com: #1C1C1E, #E51937, #767676, #D9DBD8, #FFFFFF
BG_WHITE = RGBColor(0xF5, 0xF5, 0xF3)       # Warm off-white slide background
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)         # White card backgrounds
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xDE)     # Subtle warm gray border
TEXT_PRIMARY = RGBColor(0x1C, 0x1C, 0x1E)    # SRAM near-black for headings
TEXT_BODY = RGBColor(0x3A, 0x3A, 0x3C)       # Dark gray for body text
TEXT_SECONDARY = RGBColor(0x76, 0x76, 0x76)  # SRAM mid-gray (#767676)
ACCENT_RED = RGBColor(0xE5, 0x19, 0x37)      # SRAM red - used sparingly
ACCENT_RED_SOFT = RGBColor(0xC4, 0x3B, 0x4F) # Muted red for emphasis
LIGHT_GRAY_BG = RGBColor(0xF0, 0xF0, 0xEE)  # Alternating row backgrounds
HIGHLIGHT_BG = RGBColor(0xFD, 0xF0, 0xF2)   # Very subtle red tint for highlight
GANTT_ACCENT = RGBColor(0x2D, 0x5F, 0x8A)   # Steel blue for bars and badges
GANTT_LIGHT = RGBColor(0xD4, 0xE4, 0xF0)    # Light blue for setup bars

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Shorthand
BG = BG_WHITE
CARD = CARD_BG
BORDER = CARD_BORDER
BLACK = TEXT_PRIMARY
BODY = TEXT_BODY
GRAY = TEXT_SECONDARY
RED = ACCENT_RED
FONT = "Calibri"


def set_slide_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=CARD, border=BORDER, radius=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape


def text(slide, left, top, w, h, content, size=14, color=BODY, bold=False,
         align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = align
    return box


def multitext(slide, left, top, w, h, lines, spacing=1.3):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (t, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.space_after = Pt(size * (spacing - 1) * 2)
    return box


def slide_header(slide, label, title):
    """Action-title header: small label + bold conclusion statement."""
    text(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.25),
         label, size=11, color=GRAY, bold=True)
    text(slide, Inches(0.8), Inches(0.6), Inches(10.5), Inches(0.55),
         title, size=26, color=BLACK, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.15), Inches(0.8), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def slide_footer(slide, page_num):
    """Consistent footer with page number."""
    text(slide, Inches(0.8), Inches(7.05), Inches(5), Inches(0.3),
         "SRAM AI Adoption Playbook", size=8, color=GRAY)
    text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
         str(page_num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def metric_card(slide, left, top, w, h, label, value, val_color=BLACK, sub=None):
    add_rect(slide, left, top, w, h, CARD, BORDER)
    text(slide, left + Inches(0.25), top + Inches(0.15), w - Inches(0.5), Inches(0.25),
         label, size=10, color=GRAY)
    text(slide, left + Inches(0.25), top + Inches(0.4), w - Inches(0.5), Inches(0.45),
         value, size=26, color=val_color, bold=True)
    if sub:
        text(slide, left + Inches(0.25), top + Inches(0.9), w - Inches(0.5), Inches(0.25),
             sub, size=9, color=GRAY)


def bullet_list(slide, left, top, w, h, items, color=BODY, size=12):
    lines = [(f"\u2022  {item}", size, color, False) for item in items]
    multitext(slide, left, top, w, h, lines, spacing=1.5)


def quote_box(slide, left, top, w, h, quote, attribution):
    """Styled quote callout with left accent bar."""
    add_rect(slide, left, top, w, h, CARD, BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GRAY
    bar.line.fill.background()
    text(slide, left + Inches(0.35), top + Inches(0.1), w - Inches(0.6), h - Inches(0.4),
         f"\u201C{quote}\u201D", size=12, color=BODY, italic=True)
    text(slide, left + Inches(0.35), top + h - Inches(0.3), w - Inches(0.6), Inches(0.25),
         attribution, size=10, color=GRAY, bold=True)


def arrow_right(slide, left, top, w=Inches(0.3), h=Inches(0.3)):
    """Small right-pointing triangle arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, w, h)
    shape.rotation = 90.0
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()


# --- SRAM Logo ---
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(SCRIPT_DIR, "sram_logo.png")


def add_logo_title(slide):
    """Add SRAM logo prominently on the title slide."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(1.5), Inches(1.2), width=Inches(3.0))


def add_logo(slide):
    """Add small SRAM logo to top-right corner of content slides."""
    if os.path.exists(LOGO_PATH):
        logo_w = Inches(1.2)
        slide.shapes.add_picture(
            LOGO_PATH, SLIDE_WIDTH - logo_w - Inches(0.3), Inches(0.25),
            width=logo_w)


# ============================================================
# BUILD
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank = prs.slide_layouts[6]
page = 0


# ----------------------------------------------------------
# SLIDE 1: Title
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo_title(slide)

# Thin red accent bar
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(1.5), Inches(2.4), Inches(1.2), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = RED
line.line.fill.background()

multitext(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(2.5), [
    ("SRAM LLC", 44, BLACK, True),
    ("AI Adoption Playbook", 26, GRAY, False),
    ("", 10, GRAY, False),
    ("A $10M opportunity starting with a 6-month support pilot", 16, BODY, False),
    ("", 8, GRAY, False),
    ("Prepared for Executive Leadership  |  March 2026", 13, GRAY, False),
    ("", 6, GRAY, False),
    ("Donovan Palmer, Ed Kreienberg, Will Zheng", 13, GRAY, False),
])

text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
     "AIML 950  |  Human and Machine Intelligence  |  Northwestern Kellogg",
     size=10, color=GRAY)


# ----------------------------------------------------------
# SLIDE 2: Overview of the Business
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "THE BUSINESS",
             "SRAM makes the parts that make bikes go, stop, and shift")

# Subtitle
text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
     "A $1B+ company headquartered in Chicago, powering bikes from weekend rides to the Tour de France.",
     size=14, color=GRAY)

# Product category cards
sram_cards = [
    ("Drivetrains + Braking", "SRAM",
     "Gears, chains, and disc\nbrakes for every bike"),
    ("Suspension", "RockShox",
     "Forks and shocks that\nabsorb trail impacts"),
    ("Wheels + Cockpit", "Zipp",
     "Aero wheels, handlebars,\nand stems"),
    ("Power Meters", "Quarq",
     "Sensors that measure\nrider output in watts"),
    ("GPS Bike Computer", "Hammerhead",
     "On-bike computer for\nnavigation and training data"),
    ("Pedals", "TIME",
     "Clip-in pedal systems\nfor road and mountain"),
]

n_cards = len(sram_cards)
total_avail = Inches(11.7)
intro_gap = Inches(0.25)
intro_card_w = (total_avail - intro_gap * (n_cards - 1)) / n_cards
intro_card_h = Inches(2.4)
intro_start_x = Inches(0.8)
intro_card_top = Inches(2.2)

for i, (label, brand, desc) in enumerate(sram_cards):
    cx = intro_start_x + i * (intro_card_w + intro_gap)
    add_rect(slide, int(cx), intro_card_top, int(intro_card_w), intro_card_h)
    text(slide, int(cx), intro_card_top + Inches(0.2), int(intro_card_w), Inches(0.25),
         brand, size=10, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    text(slide, int(cx), intro_card_top + Inches(0.55), int(intro_card_w), Inches(0.5),
         label, size=16, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    text(slide, int(cx) + Inches(0.15), intro_card_top + Inches(1.2),
         int(intro_card_w) - Inches(0.3), Inches(1.0),
         desc, size=11, color=BODY, align=PP_ALIGN.CENTER)

text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
     "6 brands. 1 connected ecosystem. From weekend riders to Tour de France pros.",
     size=14, color=BODY, bold=True, align=PP_ALIGN.CENTER)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 3: Competitive Advantage
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "COMPETITIVE ADVANTAGE",
             "SRAM's connected ecosystem is unmatched in cycling")

# Metrics row
metrics_data = [
    ("REVENUE", "$1B+", "Private, Chicago-founded", BLACK),
    ("BRANDS", "6", "Integrated portfolio", BLACK),
    ("MTB SHARE", "#1", "Market leader globally", BLACK),
    ("TRUSTPILOT", "1.6 / 5.0", "Support-driven churn risk", ACCENT_RED_SOFT),
]
for i, (label, value, sub, vc) in enumerate(metrics_data):
    metric_card(slide, Inches(0.8 + i * 3.05), Inches(1.6), Inches(2.8), Inches(1.25),
                label, value, vc, sub)

# Ecosystem flow
eco_items = [
    ("SRAM", "Drivetrains + Braking"),
    ("RockShox", "Suspension"),
    ("Zipp", "Wheels + Cockpit"),
    ("Quarq", "Power Meters"),
    ("Hammerhead", "GPS Bike Computer"),
    ("TIME", "Pedals"),
]
for i, (brand, desc) in enumerate(eco_items):
    left = Inches(0.8 + i * 2.05)
    card_w = Inches(1.85)
    add_rect(slide, left, Inches(3.3), card_w, Inches(0.85), CARD, BORDER)
    text(slide, left + Inches(0.15), Inches(3.38), card_w - Inches(0.3), Inches(0.25),
         brand, size=12, color=BLACK, bold=True)
    text(slide, left + Inches(0.15), Inches(3.7), card_w - Inches(0.3), Inches(0.35),
         desc, size=10, color=GRAY)
    if i < len(eco_items) - 1:
        arrow_right(slide, left + card_w + Inches(0.02), Inches(3.6))

# AXS connector bar
add_rect(slide, Inches(0.8), Inches(4.5), Inches(12.1), Inches(0.35),
         CARD, BORDER, 0.1)
text(slide, Inches(1.2), Inches(4.52), Inches(11.3), Inches(0.3),
     "AXS Wireless Ecosystem  |  All 6 brands connect wirelessly and share data",
     size=11, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

# Interview quote
quote_box(slide, Inches(0.8), Inches(5.3), Inches(12.1), Inches(0.85),
          "Shimano riders use an app to check compatibility. Our riders use an ecosystem to train, "
          "navigate, tune their suspension, and monitor their entire component stack in real time. "
          "That is a fundamentally different level of engagement.",
          "Clint Weber, VP Global Sales & Manufacturing, SRAM")

text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.25),
     "Sources: MTB share per Bicycle Retailer industry analysis; Trustpilot rating as of Feb 2026 (n=180+ reviews)",
     size=8, color=GRAY)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 4: Four-Step AI Framework Overview
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "THE STRATEGY",
             "Four steps from early wins to AI-first enterprise")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "Each step builds the infrastructure and organizational trust the next step requires.",
     size=13, color=BODY)

framework_steps = [
    ("Step 1", "Democratize AI Tooling", "Month 0-12",
     "GitHub Copilot + Claude Code\nfor 250-400 engineers\n\nMicrosoft CoPilot for 500\nbusiness function staff\n\nQuarterly AI innovation forums",
     "$400K/yr in\n~$2M/yr out  |  5x ROI"),
    ("Step 2", "Enterprise Data Consolidation", "Month 0-24",
     "AWS-native data lakehouse\non existing AWS infrastructure\n\nSAP ERP + Salesforce CRM\nconnected to central platform\n\nPrerequisite for Steps 3B + 4",
     "$6-10M investment\nUnlocks $34M+ downstream"),
    ("Step 3", "Deploy AI to Operations", "Month 0-36",
     "3A: AI support agent for\nAXS + Hammerhead dealers\n\n3B: SageMaker demand\nforecasting on pilot SKUs\n\nBounded, measurable pilots",
     "3A: ~$4.5M net Y1  |  9x\n3B: ~$4.5M net  |  9x"),
    ("Step 4", "Revenue + Innovation", "Month 36-60+",
     "AXS Intelligence Platform\n(unified rider data layer)\n\nOEM Automation + Upgrade Recs\nGenerative Design + R&D",
     "~$20M ARR at scale\n~$39M+ total at maturity"),
]

fw_card_w = Inches(2.85)
fw_gap = Inches(0.12)
fw_start = Inches(0.8)
fw_top = Inches(1.85)
fw_card_h = Inches(4.6)

for i, (step_num, title, timing, desc, fin) in enumerate(framework_steps):
    x = fw_start + i * (fw_card_w + fw_gap)
    accent = ACCENT_RED_SOFT if i == 3 else GANTT_ACCENT
    add_rect(slide, x, fw_top, fw_card_w, fw_card_h, CARD, BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, fw_top, fw_card_w, Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, x + Inches(0.2), fw_top + Inches(0.18), fw_card_w - Inches(0.4), Inches(0.22),
         step_num, size=9, color=accent, bold=True)
    text(slide, x + Inches(0.2), fw_top + Inches(0.44), fw_card_w - Inches(0.4), Inches(0.38),
         title, size=18, color=BLACK, bold=True)
    text(slide, x + Inches(0.2), fw_top + Inches(0.9), fw_card_w - Inches(0.4), Inches(0.22),
         timing, size=10, color=GRAY, bold=True)
    text(slide, x + Inches(0.2), fw_top + Inches(1.2), fw_card_w - Inches(0.4), Inches(2.2),
         desc, size=11, color=BODY)
    add_rect(slide, x + Inches(0.1), fw_top + Inches(3.35), fw_card_w - Inches(0.2), Inches(1.0),
             LIGHT_GRAY_BG, None, 0.02)
    text(slide, x + Inches(0.2), fw_top + Inches(3.45), fw_card_w - Inches(0.4), Inches(0.85),
         fin, size=10, color=BLACK, bold=True)

text(slide, Inches(0.8), Inches(6.65), Inches(11.5), Inches(0.25),
     "Step 2 runs in parallel with Steps 1 and 3A  |  "
     "Steps 3B and 4 depend on Step 2 data infrastructure",
     size=10, color=GRAY, align=PP_ALIGN.CENTER)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 5: Step 1 - Democratize AI Tooling
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 1: DEMOCRATIZE AI TOOLING",
             "Building AI fluency from the ground up is the prerequisite for every initiative")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "No new hires. No infrastructure. Deploy AI into existing tools, build organizational "
     "comfort, and identify champions before committing to Steps 3 and 4.",
     size=13, color=BODY)

step1_tools = [
    ("GitHub Copilot + Claude Code",
     "SRAM's 250-400 software engineers",
     "AI pair programming in existing IDEs.\nAccelerates the technical work Steps 2-4 require.\nIntegrates directly into existing GitHub workflows.",
     "$140K/yr",
     "75% of engineers active weekly by Month 6"),
    ("Microsoft CoPilot",
     "500 business function staff",
     "AI-assisted notes, email, and meeting\nsummarization across all business functions.\nIntegrates directly into Microsoft 365.",
     "$180K/yr",
     "60% of staff active weekly by Month 9"),
]

for i, (tool, who, desc, cost, goal) in enumerate(step1_tools):
    tx = Inches(0.8) + i * Inches(5.9)
    tc_w = Inches(5.6)
    add_rect(slide, tx, Inches(2.0), tc_w, Inches(3.5), CARD, BORDER)
    text(slide, tx + Inches(0.25), Inches(2.12), tc_w - Inches(0.5), Inches(0.3),
         tool, size=14, color=BLACK, bold=True)
    text(slide, tx + Inches(0.25), Inches(2.5), tc_w - Inches(0.5), Inches(0.22),
         who, size=11, color=GANTT_ACCENT, bold=True)
    text(slide, tx + Inches(0.25), Inches(2.8), tc_w - Inches(0.5), Inches(1.1),
         desc, size=11, color=BODY)
    text(slide, tx + Inches(0.25), Inches(3.95), tc_w - Inches(0.5), Inches(0.22),
         "Annual cost: " + cost, size=10, color=GRAY)
    text(slide, tx + Inches(0.25), Inches(4.3), tc_w - Inches(0.5), Inches(0.5),
         goal, size=10, color=BLACK, bold=True)

add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5), CARD, BORDER)
text(slide, Inches(1.2), Inches(5.68), Inches(7.5), Inches(0.32),
     "Quarterly AI Innovation Forums  |  $80K/yr  |  Surface grassroots use cases, identify AI champions",
     size=11, color=BODY)
text(slide, Inches(9.0), Inches(5.68), Inches(3.2), Inches(0.32),
     "5+ champions named by Month 12", size=10, color=BLACK, bold=True, align=PP_ALIGN.RIGHT)

add_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.65), HIGHLIGHT_BG, BORDER, 0.02)
text(slide, Inches(1.2), Inches(6.42), Inches(10.5), Inches(0.45),
     "Total investment: ~$400K/yr  |  Total benefit: ~$2.0M/yr labor-equivalent  |  "
     "Net value: ~$1.6M/yr  |  ROI: ~5x",
     size=12, color=BLACK, bold=True)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 6: Step 2 - Enterprise Data Consolidation
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 2: ENTERPRISE DATA CONSOLIDATION",
             "Data infrastructure is the prerequisite, not an optional investment")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "SRAM has no centralized ERP, CRM, or data warehouse today. Customer data lives in Shopify. "
     "AXS telemetry is siloed from Hammerhead ride data.",
     size=13, color=BODY)

arch_cols = [
    ("SOURCE SYSTEMS", [
        "SAP ERP (order + inventory)",
        "Salesforce CRM (dealers + warranty)",
        "AXS + Hammerhead telemetry",
        "Shopify (DTC orders)",
        "Quarq power data",
    ]),
    ("AWS DATA LAKEHOUSE", [
        "S3 + Glue + Redshift",
        "Unified SKU-level order data",
        "Dealer sell-through rates",
        "PowerBI + Databricks (BI layer)",
        "Data quality gates enforced",
    ]),
    ("AI / ML LAYER", [
        "Amazon Bedrock (LLM agents)",
        "Amazon SageMaker (forecasting)",
        "Amazon Kendra (doc retrieval)",
        "Salesforce AgentForce",
        "All Step 3 + 4 tools run here",
    ]),
]

arch_card_w = Inches(3.5)
arch_gap = Inches(0.45)
arch_top = Inches(2.0)
arch_h = Inches(3.0)

for i, (col_title, items) in enumerate(arch_cols):
    ax = Inches(0.8) + i * (arch_card_w + arch_gap)
    col_color = GANTT_LIGHT if i == 1 else CARD
    add_rect(slide, ax, arch_top, arch_card_w, arch_h, col_color, BORDER)
    text(slide, ax + Inches(0.2), arch_top + Inches(0.12), arch_card_w - Inches(0.4), Inches(0.25),
         col_title, size=10, color=GANTT_ACCENT, bold=True)
    for j, item in enumerate(items):
        text(slide, ax + Inches(0.25), arch_top + Inches(0.52 + j * 0.47),
             arch_card_w - Inches(0.5), Inches(0.38),
             item, size=11, color=BODY)
    if i < 2:
        arr = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                     ax + arch_card_w + Inches(0.12),
                                     arch_top + arch_h / 2 - Inches(0.1),
                                     Inches(0.2), Inches(0.2))
        arr.rotation = 90.0
        arr.fill.solid()
        arr.fill.fore_color.rgb = GRAY
        arr.line.fill.background()

metric_card(slide, Inches(0.8), Inches(5.2), Inches(2.8), Inches(1.0),
            "24-MONTH INVESTMENT", "$6-10M", BLACK, "SAP + Salesforce + Lakehouse")
metric_card(slide, Inches(3.8), Inches(5.2), Inches(2.8), Inches(1.0),
            "DIRECT ANNUAL SAVINGS", "~$1.5M/yr", BLACK, "Reporting + reconciliation")
metric_card(slide, Inches(6.8), Inches(5.2), Inches(2.8), Inches(1.0),
            "DOWNSTREAM VALUE", "$34M+/yr", ACCENT_RED_SOFT, "Steps 3B + 4 depend on this")
metric_card(slide, Inches(9.8), Inches(5.2), Inches(2.8), Inches(1.0),
            "DATA READINESS TODAY", "60-70%", GANTT_ACCENT, "AXS + Hammerhead pilot-ready")

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDES 7-10: Initiative Map (progressive build - 4 slides)
# ----------------------------------------------------------

FUNC_COLORS = {
    "SUPPORT": RGBColor(0x2D, 0x5F, 0x8A),
    "SUPPLY CHAIN": RGBColor(0x5B, 0x8C, 0x5A),
    "SALES": RGBColor(0x8B, 0x6D, 0x2E),
    "ENGINEERING": RGBColor(0x7A, 0x5C, 0x8A),
}

initiative_rows = [
    ("ENGINEERING", [
        ("Generative Design", "Physical testing is slow and expensive", False),
        ("AXS Intelligence", "Telemetry data is collected but unused", False),
        ("AI-Tuned Components", "Suspension and shifting are static", False),
    ]),
    ("SALES", [
        ("OEM Proposal Automation", "Proposal cycles are slow, manual", False),
        ("Part Recommendations", "No guided upsell at point of service", False),
        ("DTC Personalization", "Velocio has no targeting capability", False),
    ]),
    ("SUPPLY CHAIN", [
        ("Demand Forecasting", "No forecasting system exists today", False),
        ("Inventory Allocation", "Stockouts and rush shipping losses", False),
        ("Customer Analytics", "Shopify data is siloed, unused", False),
    ]),
    ("SUPPORT", [
        ("AI Support Agent", "70% of dealer questions are repetitive", True),
        ("Warranty Automation", "Fraudulent and mis-filed claims cost time", False),
        ("Compatibility Assistant", "Riders get wrong parts, return them", False),
    ]),
]

grid_left = Inches(2.4)
card_w = Inches(3.2)
card_h = Inches(0.9)

page += 1
for build_step in range(4):
    visible_up_to = build_step
    show_outline = (build_step == 3)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_logo(slide)
    slide_header(slide, "WHERE AI ADDS VALUE",
                 "12 AI initiatives across four business functions")

    if show_outline:
        sub = ("Start where domain knowledge is high and outcomes are measurable. "
               "Support is the proving ground for everything that follows.")
    else:
        sub = "AI opportunities exist across every function at SRAM."
    text(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.35),
         sub, size=13, color=BODY)

    arrow_labels = ["Quick Win", "Medium Effort", "Long-Term Bet"]
    for ci, alabel in enumerate(arrow_labels):
        ax = grid_left + Inches(ci * (3.2 + 0.15))
        text(slide, ax, Inches(1.75), card_w, Inches(0.22),
             alabel, size=9, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

    arrow_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        grid_left, Inches(1.97), Inches(10.05), Pt(2))
    arrow_line.fill.solid()
    arrow_line.fill.fore_color.rgb = BORDER
    arrow_line.line.fill.background()

    arrow_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                       grid_left + Inches(10.0), Inches(1.91),
                                       Inches(0.15), Inches(0.14))
    arrow_tri.fill.solid()
    arrow_tri.fill.fore_color.rgb = BORDER
    arrow_tri.line.fill.background()
    arrow_tri.rotation = 90.0

    for ri, (area, items) in enumerate(initiative_rows):
        row_top = Inches(2.1) + Inches(ri * (0.9 + 0.25))
        fc = FUNC_COLORS[area]

        if ri > visible_up_to:
            continue

        label_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(0.8), row_top, Pt(4), card_h)
        label_bar.fill.solid()
        label_bar.fill.fore_color.rgb = fc
        label_bar.line.fill.background()

        text(slide, Inches(0.95), row_top + Inches(0.25), Inches(1.35), Inches(0.4),
             area, size=9, color=fc, bold=True)

        for ci, (initiative, problem, highlight) in enumerate(items):
            left = grid_left + Inches(ci * (3.2 + 0.15))

            if highlight and show_outline:
                add_rect(slide, left, row_top, card_w, card_h, HIGHLIGHT_BG, RED)
                hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              left, row_top, Pt(4), card_h)
                hbar.fill.solid()
                hbar.fill.fore_color.rgb = RED
                hbar.line.fill.background()
                text(slide, left + Inches(0.2), row_top + Inches(0.08),
                     Inches(2.5), Inches(0.2),
                     "RECOMMENDED PILOT", size=8, color=RED, bold=True)
            else:
                add_rect(slide, left, row_top, card_w, card_h, CARD, BORDER)

            text(slide, left + Inches(0.2), row_top + Inches(0.25),
                 card_w - Inches(0.4), Inches(0.3),
                 initiative, size=12, color=BLACK, bold=True)
            text(slide, left + Inches(0.2), row_top + Inches(0.55),
                 card_w - Inches(0.4), Inches(0.3),
                 problem, size=9, color=GRAY)

    if show_outline:
        support_row_top = Inches(2.1) + Inches(3 * (0.9 + 0.25))
        support_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.65), support_row_top - Inches(0.08),
            Inches(11.95), card_h + Inches(0.16))
        support_box.fill.background()
        support_box.line.color.rgb = FUNC_COLORS["SUPPORT"]
        support_box.line.width = Pt(2.5)

        text(slide, Inches(10.7), support_row_top - Inches(0.32),
             Inches(1.8), Inches(0.22),
             "START HERE", size=10, color=FUNC_COLORS["SUPPORT"], bold=True,
             align=PP_ALIGN.RIGHT)

    slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 11: Step 3A - AI Customer Support
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 3A: AI CUSTOMER SUPPORT",
             "Support quality drives churn, not product quality")

text(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.35),
     "Riders love SRAM hardware. They struggle with firmware updates, "
     "part compatibility, and warranty resolution.",
     size=13, color=BODY)

# Problem metrics
metric_card(slide, Inches(0.8), Inches(1.9), Inches(3.6), Inches(1.2),
            "TRUSTPILOT RATING", "1.6 / 5.0", ACCENT_RED_SOFT,
            "Warranty friction and support delays")
metric_card(slide, Inches(4.9), Inches(1.9), Inches(3.6), Inches(1.2),
            "AUTOMATABLE", "~70%", BLACK,
            "Questions follow documented patterns")
metric_card(slide, Inches(9.0), Inches(1.9), Inches(3.6), Inches(1.2),
            "TOP COMPLAINTS", "Firmware + Pairing", ACCENT_RED_SOFT,
            "Reddit, app stores, Trustpilot")

# TODAY vs PILOT columns
today_left = Inches(0.8)
col_w = Inches(5.5)
text(slide, today_left, Inches(3.3), col_w, Inches(0.3),
     "TODAY", size=14, color=ACCENT_RED_SOFT, bold=True)

today_steps = [
    ("Dealer emails a question",
     "Agent searches a 200-page compatibility PDF manually."),
    ("Agent writes a response from scratch",
     "Types the answer, double-checks part numbers, hopes nothing changed."),
    ("Dealer waits ~4 hours for first response",
     "70% of questions follow the same patterns every time."),
]
for i, (step, detail) in enumerate(today_steps):
    y = Inches(3.7) + Inches(i * 0.85)
    add_rect(slide, today_left, y, col_w, Inches(0.75), CARD, BORDER)
    text(slide, today_left + Inches(0.2), y + Inches(0.08),
         col_w - Inches(0.4), Inches(0.25),
         step, size=11, color=BLACK, bold=True)
    text(slide, today_left + Inches(0.2), y + Inches(0.38),
         col_w - Inches(0.4), Inches(0.3),
         detail, size=10, color=GRAY)

# PILOT column
pilot_left = Inches(6.8)
text(slide, pilot_left, Inches(3.3), col_w, Inches(0.3),
     "WITH AI PILOT", size=14, color=GANTT_ACCENT, bold=True)

pilot_steps = [
    ("AI searches approved docs instantly",
     "Amazon Kendra indexes compatibility tables and firmware notes."),
    ("AI drafts a response for the agent",
     "Amazon Bedrock generates an answer. Agent reviews before sending."),
    ("Dealer gets a faster, verified answer",
     "Target first response: ~2.5 hours. Human approves every response."),
]
for i, (step, detail) in enumerate(pilot_steps):
    y = Inches(3.7) + Inches(i * 0.85)
    bg = HIGHLIGHT_BG if i == 2 else CARD
    bd = GANTT_ACCENT if i == 2 else BORDER
    add_rect(slide, pilot_left, y, col_w, Inches(0.75), bg, bd)
    text(slide, pilot_left + Inches(0.2), y + Inches(0.08),
         col_w - Inches(0.4), Inches(0.25),
         step, size=11, color=BLACK, bold=True)
    text(slide, pilot_left + Inches(0.2), y + Inches(0.38),
         col_w - Inches(0.4), Inches(0.3),
         detail, size=10, color=GRAY)

arrow_right(slide, Inches(6.35), Inches(4.5), Inches(0.35), Inches(0.35))

# Quote
quote_box(slide, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.65),
          "The first question I ask anyone skeptical about AI is: do you want to spend "
          "your day searching a 200-page compatibility PDF, or do you want to spend it "
          "talking to dealers?",
          "Clint Weber, VP Global Sales & Manufacturing, SRAM")

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 12: Step 3A - Financials and Timeline
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 3A: FINANCIALS AND TIMELINE",
             "6-month pilot: 9x ROI on $500K setup")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "Day-45 scope review, Day-90 go/no-go gate, Month 7-12 scale to full catalog.",
     size=13, color=BODY)

# LEFT: Goals and financials
text(slide, Inches(0.8), Inches(1.85), Inches(5.5), Inches(0.3),
     "SMART Goals", size=14, color=BLACK, bold=True)

smart_goals = [
    ("G1", "Cut first-response time", "-40%"),
    ("G2", "Resolve without escalation", "+15 pts"),
    ("G3", "AI draft acceptance rate", ">70%"),
    ("G4", "Agent throughput increase", "+25%"),
    ("G5", "Hold satisfaction flat", "No decline"),
]

for i, (gid, title, target) in enumerate(smart_goals):
    y = Inches(2.25) + Inches(i * 0.58)
    if i % 2 == 0:
        add_rect(slide, Inches(0.8), y - Inches(0.03), Inches(5.5), Inches(0.52),
                 LIGHT_GRAY_BG, None, 0.02)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.0), y + Inches(0.06),
                                   Inches(0.45), Inches(0.32))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GANTT_ACCENT
    badge.line.fill.background()
    badge.adjustments[0] = 0.15
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = gid
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    text(slide, Inches(1.6), y + Inches(0.04), Inches(3.0), Inches(0.3),
         title, size=12, color=BLACK, bold=True)
    text(slide, Inches(4.6), y + Inches(0.04), Inches(1.5), Inches(0.3),
         target, size=14, color=RED, bold=True)

# Financial summary below goals
metric_card(slide, Inches(0.8), Inches(5.3), Inches(2.6), Inches(0.95),
            "SETUP COST", "$500K", BLACK, "One-time integration")
metric_card(slide, Inches(3.55), Inches(5.3), Inches(2.6), Inches(0.95),
            "NET YEAR-1 VALUE", "~$4.5M", BLACK, "After operating costs")

# Kill criteria
text(slide, Inches(0.8), Inches(6.4), Inches(5.5), Inches(0.25),
     "Kill criteria: 2 weeks quality drop, any safety error, or 15% dealer opt-out",
     size=9, color=GRAY)

# RIGHT: Timeline (simplified Gantt)
text(slide, Inches(6.8), Inches(1.85), Inches(5.7), Inches(0.3),
     "6-Month Timeline", size=14, color=BLACK, bold=True)

gantt_left = Inches(8.2)
gantt_width = Inches(4.3)
month_w = gantt_width / 6

for mi, mlabel in enumerate(["M1", "M2", "M3", "M4", "M5", "M6"]):
    x = gantt_left + month_w * mi
    text(slide, x, Inches(2.2), month_w, Inches(0.2),
         mlabel, size=8, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

for di in range(1, 6):
    x = gantt_left + month_w * di
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, Inches(2.4), Pt(1), Inches(3.8))
    divider.fill.solid()
    divider.fill.fore_color.rgb = CARD_BORDER
    divider.line.fill.background()

gantt_rows = [
    ("Setup + Integration", 0, 1.0, GANTT_LIGHT),
    ("Knowledge Base Indexing", 0, 0.8, GANTT_LIGHT),
    ("AI Draft Rollout (AXS)", 1.0, 2.0, GANTT_ACCENT),
    ("AI Draft Rollout (Hammerhead)", 3.0, 2.0, GANTT_ACCENT),
    ("Response Time Tracking", 1.0, 5.0, GANTT_ACCENT),
    ("Escalation Rate Tracking", 1.0, 5.0, GANTT_ACCENT),
    ("Throughput Measurement", 2.0, 4.0, GANTT_ACCENT),
    ("CSAT Monitoring", 1.0, 5.0, GANTT_ACCENT),
    ("Weekly Quality Reviews", 1.0, 5.0, GRAY),
    ("Day-45 Checkpoint", 2.5, 0.15, RED),
    ("Day-90 Go/No-Go", 5.75, 0.15, RED),
]

row_height = Inches(0.34)
for i, (label, start, dur, color) in enumerate(gantt_rows):
    y = Inches(2.45) + row_height * i
    if i % 2 == 0:
        add_rect(slide, Inches(6.8), y, Inches(5.7), row_height,
                 LIGHT_GRAY_BG, None, 0.01)
    text(slide, Inches(6.8), y + Inches(0.04), Inches(1.35), Inches(0.25),
         label, size=7, color=BLACK, bold=True)
    bar_x = gantt_left + month_w * start
    bar_w = month_w * dur
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 int(bar_x), y + Inches(0.04),
                                 int(bar_w), Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.15

# Scope constraints row
scope_items = [
    ("SCOPE", "AXS + Hammerhead only"),
    ("GUARDRAIL", "Human approves every response"),
    ("KILL SWITCH", "Weekly quality review"),
]
for i, (label, desc) in enumerate(scope_items):
    left = Inches(6.8) + Inches(i * 1.95)
    add_rect(slide, left, Inches(6.25), Inches(1.8), Inches(0.65), CARD, BORDER)
    text(slide, left + Inches(0.1), Inches(6.28), Inches(1.6), Inches(0.18),
         label, size=7, color=GANTT_ACCENT, bold=True)
    text(slide, left + Inches(0.1), Inches(6.48), Inches(1.6), Inches(0.35),
         desc, size=8, color=BODY)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 13: Step 3B - Demand Forecasting
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 3B: DEMAND FORECASTING",
             "No forecasting system exists today; SageMaker fills that gap")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "Stockouts, markdown losses, and overproduction are addressable without touching the product. "
     "Requires the Step 2 data foundation as a hard gate.",
     size=13, color=BODY)

# TODAY vs PILOT columns
text(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.3),
     "TODAY", size=14, color=ACCENT_RED_SOFT, bold=True)

fc_steps = [
    ("1. Brand-level judgment drives production",
     "No centralized inventory visibility.\nDealer sell-through not systematically tracked."),
    ("2. Stockouts surface as surprises",
     "Rush shipping absorbs margin.\nMarkdowns clear overstock at loss."),
    ("3. Planners react to problems",
     "Decisions made without model support.\nData exists at brand level, not enterprise level."),
]
for i, (step, detail) in enumerate(fc_steps):
    y = Inches(2.3) + Inches(i * 1.1)
    add_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.95), CARD, BORDER)
    text(slide, Inches(1.0), y + Inches(0.08), Inches(5.1), Inches(0.25),
         step, size=12, color=BLACK, bold=True)
    text(slide, Inches(1.0), y + Inches(0.38), Inches(5.1), Inches(0.5),
         detail, size=10, color=GRAY)

text(slide, Inches(6.8), Inches(1.9), Inches(5.7), Inches(0.3),
     "WITH SAGEMAKER PILOT", size=14, color=GANTT_ACCENT, bold=True)

pilot_steps_3b = [
    ("Month 22-27: Model development",
     "SageMaker trained on unified lakehouse data.\nBack-test validates MAPE target of 15% or less."),
    ("Month 27-30: Shadow mode",
     "Model runs in parallel. Planners receive\nrecommendations but make decisions as normal."),
    ("Month 30-36: Active pilot",
     "Planners act on model recommendations\nwith full override authority. Overrides logged."),
]
for i, (step, detail) in enumerate(pilot_steps_3b):
    y = Inches(2.3) + Inches(i * 1.1)
    bg = HIGHLIGHT_BG if i == 2 else CARD
    bd = GANTT_ACCENT if i == 2 else BORDER
    add_rect(slide, Inches(6.8), y, Inches(5.7), Inches(0.95), bg, bd)
    text(slide, Inches(7.0), y + Inches(0.08), Inches(5.3), Inches(0.25),
         step, size=12, color=BLACK, bold=True)
    text(slide, Inches(7.0), y + Inches(0.38), Inches(5.3), Inches(0.5),
         detail, size=10, color=GRAY)

arrow_right(slide, Inches(6.35), Inches(3.3), Inches(0.35), Inches(0.35))

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 14: Step 3B - Financials and Timeline
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 3B: FINANCIALS AND TIMELINE",
             "Demand forecasting delivers 9x ROI once data infrastructure is live")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "Begins at Month 22 after the Step 2 data lakehouse is operational. "
     "Shadow mode validates accuracy before planners act on recommendations.",
     size=13, color=BODY)

# Financial metric cards
metric_card(slide, Inches(0.8), Inches(2.0), Inches(2.8), Inches(1.25),
            "SETUP INVESTMENT", "~$500K", BLACK, "SageMaker + data science lead")
metric_card(slide, Inches(3.8), Inches(2.0), Inches(2.8), Inches(1.25),
            "ANNUAL GROSS BENEFIT", "$5.2M", BLACK, "Stockouts + markdowns reduction")
metric_card(slide, Inches(6.8), Inches(2.0), Inches(2.8), Inches(1.25),
            "NET ANNUAL VALUE", "~$4.5M", BLACK, "After $700K operating cost")
metric_card(slide, Inches(9.8), Inches(2.0), Inches(2.8), Inches(1.25),
            "ONE-TIME BONUS", "$15-20M", GANTT_ACCENT, "Working capital release from inventory")

# Success criteria
text(slide, Inches(0.8), Inches(3.55), Inches(5.5), Inches(0.3),
     "Success Criteria", size=14, color=BLACK, bold=True)

fc_goals = [
    ("Forecast accuracy (MAPE)", "15% or less", "Back-tested before launch"),
    ("Stockout reduction", "-30%", "Measured against 12-month baseline"),
    ("Markdown reduction", "-20%", "Tracked at SKU level"),
    ("Planner adoption", ">80%", "Recommendations reviewed weekly"),
]
for i, (title, target, detail) in enumerate(fc_goals):
    y = Inches(4.0) + Inches(i * 0.55)
    if i % 2 == 0:
        add_rect(slide, Inches(0.8), y - Inches(0.03), Inches(5.5), Inches(0.5),
                 LIGHT_GRAY_BG, None, 0.02)
    text(slide, Inches(1.0), y + Inches(0.04), Inches(2.5), Inches(0.3),
         title, size=11, color=BLACK, bold=True)
    text(slide, Inches(3.5), y + Inches(0.02), Inches(1.0), Inches(0.3),
         target, size=13, color=RED, bold=True)
    text(slide, Inches(4.6), y + Inches(0.06), Inches(1.7), Inches(0.25),
         detail, size=9, color=GRAY)

# Timeline
text(slide, Inches(6.8), Inches(3.55), Inches(5.7), Inches(0.3),
     "Timeline", size=14, color=BLACK, bold=True)

fc_timeline = [
    ("Month 0-24", "Data foundation (Step 2)", GANTT_LIGHT,
     "SAP, Salesforce, and lakehouse build.\nHard gate for forecasting launch."),
    ("Month 22-27", "Model development", GANTT_ACCENT,
     "SageMaker trained on unified data.\nBack-test validates accuracy."),
    ("Month 27-30", "Shadow mode", GANTT_ACCENT,
     "Planners compare model vs. judgment.\nNo operational changes yet."),
    ("Month 30-36", "Active pilot", RED,
     "Live recommendations with override.\nOverrides logged for model retraining."),
]
for i, (timing, label, color, desc) in enumerate(fc_timeline):
    y = Inches(4.0) + Inches(i * 0.75)
    add_rect(slide, Inches(6.8), y, Inches(5.7), Inches(0.65), CARD, BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), y, Pt(4), Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    text(slide, Inches(7.15), y + Inches(0.04), Inches(1.3), Inches(0.2),
         timing, size=9, color=color, bold=True)
    text(slide, Inches(8.5), y + Inches(0.04), Inches(3.8), Inches(0.2),
         label, size=11, color=BLACK, bold=True)
    text(slide, Inches(7.15), y + Inches(0.28), Inches(5.1), Inches(0.3),
         desc, size=9, color=GRAY)

# ROI summary
add_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.55), HIGHLIGHT_BG, BORDER, 0.02)
text(slide, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.35),
     "ROI: ~9x  |  Net annual value: ~$4.5M  |  Working capital release: $15-20M (one-time)  |  "
     "Payback period: <6 months after go-live",
     size=12, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDES 15-16: What Could Be (2031 and beyond)
# ----------------------------------------------------------
page += 1

visions = [
    ("AXS Intelligence Platform",
     "Unified rider dashboard: power, gearing, suspension, heart rate "
     "in one view. No competitor has the product breadth to build this.",
     "NEW REVENUE: Premium subscription tier"),
    ("Predictive Maintenance",
     "Connected components predict chain and brake pad replacement "
     "before failure. Installed base becomes recurring revenue.",
     "NEW REVENUE: Per-device monitoring subscription"),
    ("AI-Tuned Performance",
     "Suspension auto-adjusts to terrain. Shifting optimizes for rider "
     "power. Products improve continuously after purchase.",
     "NEW REVENUE: Premium optimization tier"),
    ("Dealer Intelligence",
     "Telemetry tells dealers which parts approach end-of-life in their "
     "area. Proactive ordering replaces reactive.",
     "NEW REVENUE: Dealer analytics license"),
]

for build in range(2):
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_logo(slide)
    slide_header(slide, "WHAT COULD BE",
                 "Hardware company to performance intelligence by 2031")

    text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
         "Each capability creates a new recurring revenue stream. "
         "SRAM's installed base of connected components becomes a platform, not a one-time sale.",
         size=13, color=BODY)

    # AXS Intelligence Platform - always shown
    axs_top = Inches(2.0)
    axs_w = Inches(11.7)
    axs_h = Inches(1.5)
    add_rect(slide, Inches(0.8), axs_top, axs_w, axs_h, HIGHLIGHT_BG, RED)
    hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), axs_top, Pt(4), axs_h)
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = RED
    hbar.line.fill.background()
    text(slide, Inches(1.2), axs_top + Inches(0.12), Inches(4.0), Inches(0.2),
         "THE UNLOCK", size=10, color=RED, bold=True)
    text(slide, Inches(1.2), axs_top + Inches(0.4), Inches(5.0), Inches(0.4),
         visions[0][0], size=20, color=BLACK, bold=True)
    text(slide, Inches(1.2), axs_top + Inches(0.9), Inches(10.5), Inches(0.5),
         visions[0][1], size=13, color=BODY)
    text(slide, Inches(8.0), axs_top + Inches(0.4), Inches(4.2), Inches(0.3),
         visions[0][2], size=11, color=RED, bold=True, align=PP_ALIGN.RIGHT)

    if build == 1:
        for i, (title, desc, rev) in enumerate(visions[1:]):
            left = Inches(0.8) + Inches(i * 4.0)
            top = Inches(4.0)
            vc_w = Inches(3.75)
            vc_h = Inches(1.5)
            add_rect(slide, left, top, vc_w, vc_h, CARD, BORDER)
            text(slide, left + Inches(0.25), top + Inches(0.15),
                 vc_w - Inches(0.5), Inches(0.3),
                 title, size=13, color=BLACK, bold=True)
            text(slide, left + Inches(0.25), top + Inches(0.5),
                 vc_w - Inches(0.5), Inches(0.55),
                 desc, size=11, color=BODY)
            text(slide, left + Inches(0.25), top + Inches(1.15),
                 vc_w - Inches(0.5), Inches(0.25),
                 rev, size=9, color=RED, bold=True)

        # Flywheel
        add_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.45),
                 CARD, BORDER, 0.1)
        red_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.8), Inches(6.0), Pt(4), Inches(0.45))
        red_bar.fill.solid()
        red_bar.fill.fore_color.rgb = RED
        red_bar.line.fill.background()
        text(slide, Inches(1.2), Inches(6.04), Inches(11.0), Inches(0.35),
             "Hardware sells data access  \u2192  Data improves performance  "
             "\u2192  Performance sells hardware",
             size=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

    slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 17: Step 4 - New Revenue Streams
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "STEP 4: NEW REVENUE STREAMS",
             "Premium AXS, subscription model, and four distinct revenue paths")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "Each initiative is viable only after Step 2 data infrastructure is live. "
     "Combined, they represent $39M+/yr at maturity.",
     size=13, color=BODY)

step4_cards = [
    ("AXS Intelligence Platform",
     "Unified rider data: AXS drivetrain, RockShox,\nQuarq power, Hammerhead GPS in one view.\nNo competitor has the breadth to build this.",
     "$3-5M build", "~$20M ARR at scale", "5-7x"),
    ("Generative Design + R&D",
     "AI-assisted design, digital twin simulation,\nautomated patent search. Field telemetry trains\nAI on actual performance data at scale.",
     "$2-4M build", "$3.75-14M/yr benefit", "3-5x"),
    ("OEM Proposal Automation",
     "Salesforce AgentForce + CPQ accelerates\nspec proposals and automates documentation.\n1% OEM win rate improvement = $2.5M/yr.",
     "$500K-1M build", "~$3.0M/yr benefit", "4-5x"),
    ("Compatible Upgrade Recs",
     "SageMaker recommendation engine lifts basket\nvalue through dealer and DTC channels.\nPredictive maintenance drives repeat purchases.",
     "$500K-800K build", "$2.5-5M/yr benefit", "4-7x"),
]

s4_card_w = Inches(5.7)
s4_card_h = Inches(2.1)
s4_gap_x = Inches(0.2)
s4_gap_y = Inches(0.2)
s4_start_x = Inches(0.8)
s4_start_y = Inches(1.85)

for i, (title, desc, invest, benefit, roi) in enumerate(step4_cards):
    row = i // 2
    col = i % 2
    cx = s4_start_x + col * (s4_card_w + s4_gap_x)
    cy = s4_start_y + row * (s4_card_h + s4_gap_y)
    add_rect(slide, cx, cy, s4_card_w, s4_card_h, CARD, BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, s4_card_w, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    text(slide, cx + Inches(0.2), cy + Inches(0.12), s4_card_w - Inches(0.4), Inches(0.3),
         title, size=13, color=BLACK, bold=True)
    text(slide, cx + Inches(0.2), cy + Inches(0.5), s4_card_w - Inches(0.4), Inches(0.9),
         desc, size=11, color=BODY)
    text(slide, cx + Inches(0.2), cy + Inches(1.5), Inches(1.8), Inches(0.35),
         invest, size=10, color=GRAY)
    text(slide, cx + Inches(2.1), cy + Inches(1.5), Inches(2.2), Inches(0.35),
         benefit, size=10, color=BLACK, bold=True)
    text(slide, cx + Inches(4.4), cy + Inches(1.5), Inches(1.1), Inches(0.35),
         roi + " ROI", size=10, color=RED, bold=True, align=PP_ALIGN.RIGHT)

add_rect(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.55), CARD, BORDER, 0.1)
red_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(6.35), Pt(4), Inches(0.55))
red_b.fill.solid()
red_b.fill.fore_color.rgb = RED
red_b.line.fill.background()
text(slide, Inches(1.2), Inches(6.41), Inches(11.0), Inches(0.42),
     "Step 4 total: ~$6-11M investment  |  $28-42M/yr at scale  |  "
     "Full program (Steps 1-4): ~$15-22M in  |  $39M+/yr at maturity  |  ~4x blended ROI",
     size=12, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 18: What It Takes
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "WHAT IT TAKES",
             "Talent, data, and leadership required at each phase")

enabler_cols = [
    ("TALENT", [
        ("Steps 1 + 3A", "1 integration engineer + 1 product owner\n"
         "for the support pilot. GitHub Copilot for\n"
         "engineers; CoPilot for business staff.\n"
         "CTO + COO name program sponsors."),
        ("Step 2+", "3-5 data engineers + 1 senior data architect\n"
         "for the lakehouse. Add ML lead at Month 20\n"
         "for demand forecasting. Constraint is data\n"
         "engineering capacity, not AI tooling."),
    ]),
    ("DATA", [
        ("Today", "AXS + Hammerhead data is pilot-ready.\n"
         "Weber: 60-70% of the way to AI-ready\n"
         "in this domain. Everything else is siloed\n"
         "by brand. Shopify is the biggest gap."),
        ("Needed", "Unified lakehouse linking SAP, Salesforce,\n"
         "telemetry, and Shopify. Hard gate before\n"
         "Step 3B and Step 4 launch.\n"
         "24-month build; start now."),
    ]),
    ("LEADERSHIP", [
        ("Before pilot", "CEO names a dedicated C-suite AI Sponsor\n"
         "with primary accountability. Not a committee.\n"
         "Not a secondary responsibility.\n"
         "Highest-leverage decision in the plan."),
        ("Phase gate", "Sponsor resolves cross-functional conflicts,\n"
         "enforces data centralization, and holds\n"
         "sequencing discipline across all four steps.\n"
         "No sponsor, no Step 2."),
    ]),
]

col_w = Inches(3.75)
col_gap = Inches(0.2)
col_start = Inches(0.8)

for ci, (col_title, items) in enumerate(enabler_cols):
    cx = col_start + ci * (col_w + col_gap)
    text(slide, cx, Inches(1.5), col_w, Inches(0.3),
         col_title, size=14, color=BLACK, bold=True)
    for ri, (phase, detail) in enumerate(items):
        ry = Inches(2.0) + Inches(ri * 2.2)
        card_ht = Inches(2.0)
        add_rect(slide, cx, ry, col_w, card_ht, CARD, BORDER)
        text(slide, cx + Inches(0.2), ry + Inches(0.12), col_w - Inches(0.4), Inches(0.25),
             phase.upper(), size=9, color=GRAY, bold=True)
        text(slide, cx + Inches(0.2), ry + Inches(0.4), col_w - Inches(0.4), Inches(1.4),
             detail, size=11, color=BODY)

quote_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
          "We are 60 to 70 percent of the way to AI-ready data "
          "in the AXS and Hammerhead domain.",
          "Clint Weber, VP Global Sales & Manufacturing, SRAM")

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 19: Trade-offs and Risks
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "TRADE-OFFS AND RISKS",
             "Three risks worth naming before the decision")

tradeoffs = [
    ("Dealer trust is fragile",
     "One confident wrong answer from the AI reaches a dealer, and the support "
     "team loses credibility it spent years building. Human-in-the-loop and "
     "weekly quality reviews mitigate this, but they also slow throughput gains.",
     "MITIGATION: Kill criteria trigger rollback at first safety error"),
    ("Data consolidation has no executive sponsor",
     "Weber owns sales and manufacturing but not the CRM, ERP, or inventory systems. "
     "Phase 2 requires a unified data layer across those systems. Without a "
     "senior sponsor, data integration stalls and the pilot stays isolated.",
     "MITIGATION: CEO assigns cross-functional data owner in parallel"),
    ("AI hype creates unrealistic internal expectations",
     "SRAM's engineering culture values precision. If leadership frames AI as a "
     "silver bullet rather than a tool, early imperfections will fuel skepticism "
     "and internal resistance to future phases.",
     "MITIGATION: Frame pilot as experiment with defined success and failure criteria"),
]

for i, (title, desc, mitigation) in enumerate(tradeoffs):
    top = Inches(1.5) + Inches(i * 1.75)
    add_rect(slide, Inches(0.8), top, Inches(11.5), Inches(1.5), CARD, BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.8), top, Pt(4), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_RED_SOFT
    bar.line.fill.background()
    text(slide, Inches(1.2), top + Inches(0.1), Inches(10.5), Inches(0.3),
         title, size=14, color=BLACK, bold=True)
    text(slide, Inches(1.2), top + Inches(0.45), Inches(10.5), Inches(0.6),
         desc, size=11, color=BODY)
    text(slide, Inches(1.2), top + Inches(1.1), Inches(10.5), Inches(0.25),
         mitigation, size=10, color=GRAY, bold=True)

slide_footer(slide, page)


# ----------------------------------------------------------
# SLIDE 20: Clear Next Steps for Leadership
# ----------------------------------------------------------
page += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "NEXT STEPS",
             "Three actions for SRAM's CEO")

actions = [
    ("1", "Launch the 6-Month Support Program",
     "AXS and Hammerhead dealer support. One product owner, one integration "
     "engineer. Human approval on all outputs. Day-45 scope review, Day-90 "
     "go/no-go gate. ~$4.5M net Year-1 value at 9x ROI."),
    ("2", "Name the C-Suite AI Sponsor",
     "One person. Primary accountability. Not a committee, not a side "
     "responsibility. This single decision is the highest-leverage action in the "
     "plan. The program does not survive Step 2 organizational friction without it."),
    ("3", "Authorize the Step 2 Data Investment",
     "Begin SAP and Salesforce vendor selection now, in parallel with the pilot. "
     "Every month of delay is a month of delay on demand forecasting, AXS "
     "Intelligence, and OEM automation. $6-10M unlocks $34M+/yr at maturity."),
]

for i, (num, title, desc) in enumerate(actions):
    top = Inches(1.5) + Inches(i * 1.55)
    add_rect(slide, Inches(0.8), top, Inches(11.5), Inches(1.3), CARD, BORDER)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), top + Inches(0.25),
                                    Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLACK
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = BG_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    text(slide, Inches(2.0), top + Inches(0.1), Inches(9.8), Inches(0.35),
         title, size=18, color=BLACK, bold=True)
    text(slide, Inches(2.0), top + Inches(0.5), Inches(9.8), Inches(0.6),
         desc, size=12, color=BODY)

# Bottom bar
add_rect(slide, Inches(0.8), Inches(6.25), Inches(11.5), Inches(0.6), CARD, BORDER, 0.1)
red_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.8), Inches(6.25), Pt(4), Inches(0.6))
red_bar.fill.solid()
red_bar.fill.fore_color.rgb = RED
red_bar.line.fill.background()
multitext(slide, Inches(1.3), Inches(6.28), Inches(10.5), Inches(0.55), [
    ("Support pilot: 9x ROI on $500K setup  |  Full program: $39M+/yr at maturity on $15-22M investment", 13, BLACK, True),
    ("Downside bounded by pilot kill criteria  |  Upside scales with SRAM's connected data advantage", 11, GRAY, False),
], 1.1)

slide_footer(slide, page)


# ===========================================================
# APPENDIX SLIDES
# ===========================================================

# ----------------------------------------------------------
# APPENDIX A1: Full Financial Breakdown
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "APPENDIX",
             "Financial detail across three scenarios")

headers = ["Category", "Conservative", "Expected", "Upside"]
rows = [
    ("Support assistant savings", "$0.5M", "$1.6M", "$4.3M"),
    ("Forecasting + inventory", "$0.7M", "$1.6M", "$4.5M"),
    ("Documentation + translation", "$0.1M", "$0.24M", "$0.9M"),
    ("Customer retention revenue", "$1.5M", "$3.0M", "$4.5M"),
    ("Package size increase", "$2.0M", "$4.0M", "$8.0M"),
    ("Proposal win-rate revenue", "$1.3M", "$2.5M", "$3.8M"),
]
totals = [
    ("Gross value", "$6.1M", "$12.9M", "$26.0M"),
    ("Setup costs", "($1.2M)", "($2.0M)", "($2.5M)"),
    ("Operating costs", "($0.4M)", "($0.7M)", "($1.0M)"),
    ("Net Year-1 value", "$4.5M", "$10.2M", "$22.5M"),
    ("Return on spend", "2.8x", "3.8x", "6.4x"),
]

for j, h in enumerate(headers):
    x = Inches(0.8 + j * 3.0) if j > 0 else Inches(0.8)
    text(slide, x, Inches(1.5), Inches(3.0), Inches(0.3),
         h, size=11, color=GRAY, bold=True)

for i, (cat, c, e, u) in enumerate(rows):
    y = Inches(1.9 + i * 0.4)
    if i % 2 == 0:
        add_rect(slide, Inches(0.8), y - Inches(0.02), Inches(11.5), Inches(0.38),
                 LIGHT_GRAY_BG, None, 0.02)
    text(slide, Inches(0.8), y, Inches(3.0), Inches(0.3), cat, size=11, color=BODY)
    text(slide, Inches(3.8), y, Inches(3.0), Inches(0.3), c, size=11, color=BODY)
    text(slide, Inches(6.8), y, Inches(3.0), Inches(0.3), e,
         size=11, color=BLACK, bold=True)
    text(slide, Inches(9.8), y, Inches(3.0), Inches(0.3), u, size=11, color=BODY)

sep_y = Inches(1.9 + len(rows) * 0.4)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), sep_y, Inches(11.5), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = BORDER
line.line.fill.background()

for i, (cat, c, e, u) in enumerate(totals):
    y = sep_y + Inches(0.15 + i * 0.4)
    is_key = cat in ("Net Year-1 value", "Return on spend")
    text(slide, Inches(0.8), y, Inches(3.0), Inches(0.3), cat,
         size=11, color=BLACK, bold=is_key)
    text(slide, Inches(3.8), y, Inches(3.0), Inches(0.3), c,
         size=11, color=BLACK if is_key else BODY, bold=is_key)
    text(slide, Inches(6.8), y, Inches(3.0), Inches(0.3), e,
         size=11, color=BLACK, bold=True)
    text(slide, Inches(9.8), y, Inches(3.0), Inches(0.3), u,
         size=11, color=BLACK if is_key else BODY, bold=is_key)

text(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
     "APPENDIX A1", size=8, color=GRAY)


# ----------------------------------------------------------
# APPENDIX A2: Interview Highlights
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "APPENDIX",
             "Simulated Interview with Clint Weber, VP Global Sales & Manufacturing")

text(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.3),
     "AI-simulated interview  |  March 2026  |  Key quotes and insights",
     size=12, color=GRAY)

quotes = [
    ("On data readiness",
     "The areas where we've under-invested in process, data, and cross-functional "
     "integration are exactly the areas where AI could close the gap fastest. "
     "Customer support. Demand forecasting. Compatibility and documentation management.",
     "Validates bounded pilot scope"),
    ("On cultural barriers",
     "We built a Formula 1 car and then had a regional pit crew. The support "
     "infrastructure has lagged behind the hardware. That is a cultural artifact "
     "as much as a resource one.",
     "Human-in-the-loop is mandatory"),
    ("On competitive advantage",
     "Shimano riders use an app to check compatibility. Our riders use an "
     "ecosystem to train, navigate, tune their suspension, and monitor their "
     "entire component stack in real time. That rider does not churn.",
     "Hammerhead is SRAM's real AI moat"),
    ("On implementation discipline",
     "If it doesn't hit the metric, you shut it down, you learn from it, and "
     "you redeploy the resources. That discipline is what separates companies "
     "that build real AI capability from companies that accumulate AI debt.",
     "Aligns with our 90-day pilot targets"),
]

for i, (topic, quote, implication) in enumerate(quotes):
    y = Inches(1.85 + i * 1.3)
    add_rect(slide, Inches(0.8), y, Inches(11.8), Inches(1.1), CARD, BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.8), y, Pt(4), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GRAY
    bar.line.fill.background()
    text(slide, Inches(1.2), y + Inches(0.05), Inches(2.0), Inches(0.2),
         topic.upper(), size=9, color=GRAY, bold=True)
    text(slide, Inches(1.2), y + Inches(0.28), Inches(8.0), Inches(0.6),
         f"\u201C{quote}\u201D", size=11, color=BODY, italic=True)
    text(slide, Inches(9.5), y + Inches(0.35), Inches(2.8), Inches(0.5),
         f"\u2192 {implication}", size=10, color=BLACK, bold=True)

text(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
     "APPENDIX A2", size=8, color=GRAY)


# ----------------------------------------------------------
# APPENDIX A3: 60-Month Roadmap
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "APPENDIX",
             "Five years from pilot to AI-first enterprise")

text(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
     "Steps 1 and 3A start immediately. Step 2 runs in parallel. "
     "Steps 3B and 4 depend on Step 2 data infrastructure.",
     size=13, color=BODY)

rm_left = Inches(3.85)
rm_width = Inches(8.7)
half_w = rm_width / 10
half_labels = ["H1 Y1", "H2 Y1", "H1 Y2", "H2 Y2", "H1 Y3",
               "H2 Y3", "H1 Y4", "H2 Y4", "H1 Y5", "H2 Y5"]

for yr in range(5):
    x = rm_left + yr * (half_w * 2)
    text(slide, x, Inches(1.65), half_w * 2, Inches(0.2),
         f"Year {yr + 1}", size=9, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

for ci, label in enumerate(half_labels):
    x = rm_left + ci * half_w
    text(slide, x, Inches(1.88), half_w, Inches(0.18),
         label, size=7, color=GRAY, align=PP_ALIGN.CENTER)

for yr in range(1, 5):
    x = rm_left + yr * (half_w * 2)
    dv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.65), Pt(1), Inches(4.75))
    dv.fill.solid()
    dv.fill.fore_color.rgb = CARD_BORDER
    dv.line.fill.background()

roadmap_rows = [
    ("AI Productivity Tools", "STEP 1", 0, 4, GANTT_LIGHT, "5x ROI"),
    ("SAP ERP Implementation", "STEP 2", 0, 3, GANTT_ACCENT, ""),
    ("Salesforce + AWS Lakehouse", "STEP 2", 0, 4, GANTT_ACCENT, ""),
    ("Warranty Support Pilot", "STEP 3A", 0, 2, FUNC_COLORS["SUPPORT"], "9x ROI"),
    ("Scale to Full Catalog", "STEP 3A", 2, 4, FUNC_COLORS["SUPPORT"], ""),
    ("Demand Forecasting", "STEP 3B", 4, 4, FUNC_COLORS["SUPPLY CHAIN"], "9x ROI"),
    ("OEM Proposal Automation", "STEP 4", 6, 2, RED, "4-5x"),
    ("AXS Intelligence Platform", "STEP 4", 6, 4, RED, "$20M ARR"),
    ("Upgrade Recommendations", "STEP 4", 7, 3, RED, "5-7x"),
    ("Generative Design / R&D", "STEP 4", 7, 3, ACCENT_RED_SOFT, "3-5x"),
]

rm_row_h = Inches(0.37)
prev_step_rm = None
for i, (label, step, start, dur, color, fin_note) in enumerate(roadmap_rows):
    y = Inches(2.1) + rm_row_h * i
    if i % 2 == 0:
        add_rect(slide, Inches(0.8), y, Inches(11.7), rm_row_h, LIGHT_GRAY_BG, None, 0.01)
    if step != prev_step_rm:
        fc_rm = (GANTT_ACCENT if step in ("STEP 1", "STEP 2") else
                 FUNC_COLORS["SUPPORT"] if step == "STEP 3A" else
                 FUNC_COLORS["SUPPLY CHAIN"] if step == "STEP 3B" else RED)
        lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Pt(3), rm_row_h)
        lb.fill.solid()
        lb.fill.fore_color.rgb = fc_rm
        lb.line.fill.background()
        text(slide, Inches(0.95), y + Inches(0.06), Inches(0.8), Inches(0.25),
             step, size=8, color=fc_rm, bold=True)
        prev_step_rm = step
    text(slide, Inches(1.8), y + Inches(0.06), Inches(1.9), Inches(0.25),
         label, size=9, color=BLACK)
    bar_x = rm_left + start * half_w
    bar_w = dur * half_w
    bar_s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   int(bar_x), y + Inches(0.06),
                                   int(bar_w), Inches(0.22))
    bar_s.fill.solid()
    bar_s.fill.fore_color.rgb = color
    bar_s.line.fill.background()
    bar_s.adjustments[0] = 0.15
    if fin_note and dur >= 2:
        text(slide, int(bar_x) + Inches(0.06), y + Inches(0.06),
             int(bar_w) - Inches(0.12), Inches(0.22),
             fin_note, size=7,
             color=RGBColor(0xFF, 0xFF, 0xFF) if color != GANTT_LIGHT else BLACK)

text(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.25),
     "Total program: ~$15-22M over 5 years  |  "
     "Expected annual value at maturity: $39M+/yr  |  Blended ROI: ~4x",
     size=10, color=GRAY, align=PP_ALIGN.CENTER)

text(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
     "APPENDIX A3", size=8, color=GRAY)


# ----------------------------------------------------------
# APPENDIX A4: Phasing and Organizational Requirements
# ----------------------------------------------------------
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_logo(slide)
slide_header(slide, "APPENDIX",
             "Phasing, talent, and organizational enablers")

phase_details = [
    ("Phase 1: Support Pilot", "0-90 days",
     ["1 integration engineer (internal or contractor)",
      "1 product owner from support organization",
      "Zendesk AI + Amazon Kendra + Bedrock",
      "No new AI/ML hires required"]),
    ("Phase 2: Data Foundation", "90-180 days",
     ["Add product managers and QA capacity",
      "Unified data layer linking CRM + inventory + telemetry",
      "AWS SageMaker for forecasting",
      "ERP system evaluation"]),
    ("Phase 3: Revenue Acceleration", "180-365 days",
     ["Salesforce Einstein/Agentforce for OEM proposals",
      "Compatibility rules engine for recommendations",
      "Cross-functional product + sales alignment",
      "Expected $2.5M from 1% win-rate improvement"]),
    ("Phase 4: AI-First Culture", "Year 2+",
     ["AXS Intelligence Platform development",
      "Digital twin and generative design capability",
      "Predictive maintenance subscription model",
      "AI-adjusted component performance"]),
]

for i, (title, timing, items) in enumerate(phase_details):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 2.7)
    card_w = Inches(5.95)
    card_h = Inches(2.45)

    bg_color = HIGHLIGHT_BG if i == 0 else CARD
    border_color = RED if i == 0 else BORDER
    add_rect(slide, left, top, card_w, card_h, bg_color, border_color)
    text(slide, left + Inches(0.25), top + Inches(0.1), card_w - Inches(0.5), Inches(0.3),
         title, size=14, color=BLACK, bold=True)
    text(slide, left + Inches(0.25), top + Inches(0.35), Inches(2.0), Inches(0.2),
         timing, size=10, color=RED if i == 0 else GRAY, bold=True)
    bullet_list(slide, left + Inches(0.25), top + Inches(0.65), card_w - Inches(0.5),
                Inches(1.6), items, BODY, 11)

text(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
     "APPENDIX A4", size=8, color=GRAY)


# ----------------------------------------------------------
# SAVE
# ----------------------------------------------------------
script_dir = os.path.dirname(os.path.abspath(__file__))
repo_root = os.path.dirname(script_dir)
out_dir = os.path.join(repo_root, "deliverables")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "SRAM_AI_Adoption_Playbook.pptx")
prs.save(out_path)
print(f"Saved presentation to {out_path}")
n_slides = len(prs.slides)
n_appendix = 4
print(f"Slides: {n_slides} ({n_slides - n_appendix} content + {n_appendix} appendix)")
