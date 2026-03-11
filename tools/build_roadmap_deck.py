#!/usr/bin/env python3
"""
Build the SRAM AI Adoption Strategy Roadmap presentation.

McKinsey-style Crawl-Walk-Run framework with visual storytelling.
CEO audience. Clean, high-contrast design with strategic callouts.

Usage:
    python3 tools/build_roadmap_deck.py

Output:
    deliverables/AI_Adoption_Strategy_Roadmap.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- McKinsey-inspired palette ---
BG = RGBColor(0xFF, 0xFF, 0xFF)             # White background
BG_DARK = RGBColor(0x1B, 0x2A, 0x4A)        # Dark navy for title/section slides
BG_LIGHT = RGBColor(0xF7, 0xF8, 0xFA)       # Very light gray
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # White cards
BLACK = RGBColor(0x1A, 0x1A, 0x2E)          # Near-black text
BODY = RGBColor(0x3A, 0x3A, 0x4A)           # Dark gray body text
GRAY = RGBColor(0x8B, 0x8B, 0x9A)           # Muted gray
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)     # Borders, dividers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Phase colors
CRAWL_COLOR = RGBColor(0x22, 0x8B, 0x22)    # Forest green
WALK_COLOR = RGBColor(0x2563, 0xEB, 0x00)[0:3] if False else RGBColor(0x25, 0x63, 0xEB)  # Blue
RUN_COLOR = RGBColor(0xE5, 0x19, 0x37)      # SRAM red

# Accent
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)    # McKinsey blue accent
ACCENT_TEAL = RGBColor(0x06, 0xB6, 0xD4)    # Teal for highlights

FONT = "Calibri"
FONT_SERIF = "Georgia"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# === HELPERS ===

def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill=CARD_WHITE, border=None, radius=0.04, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        shape.adjustments[0] = radius
    return shape


def add_line(slide, x1, y1, x2, y2, color=LIGHT_GRAY, width=1):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def txt(slide, left, top, w, h, content, size=14, color=BODY, bold=False,
        align=PP_ALIGN.LEFT, italic=False, font=FONT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    return box


def multi(slide, left, top, w, h, lines, line_spacing=1.15):
    """lines = list of (text, size, color, bold, [font])"""
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        t, size, color, bold = line[0], line[1], line[2], line[3]
        font_name = line[4] if len(line) > 4 else FONT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(size * 0.4)
        p.line_spacing = Pt(size * line_spacing)
    return box


def bullet_list(slide, left, top, w, h, items, size=13, color=BODY, bullet_char="\u2022"):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(6)
        p.line_spacing = Pt(size * 1.3)
    return box


def phase_badge(slide, left, top, label, color, w=Inches(1.4), h=Inches(0.38)):
    shape = add_shape(slide, left, top, w, h, fill=color, border=None, radius=0.15)
    shape.text_frame.word_wrap = False
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    return shape


# === SLIDE BUILDERS ===

def slide_title(prs):
    """Dark navy title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide, BG_DARK)

    # Accent line
    add_shape(slide, Inches(1.2), Inches(2.6), Inches(0.8), Pt(4), fill=ACCENT_BLUE, border=None, radius=0.5)

    txt(slide, Inches(1.2), Inches(2.85), Inches(11), Inches(1.2),
        "AI Adoption Strategy Roadmap", size=40, color=WHITE, bold=True)

    txt(slide, Inches(1.2), Inches(4.1), Inches(9), Inches(0.6),
        "A phased approach to building AI capability at SRAM", size=20, color=GRAY, italic=True, font=FONT_SERIF)

    txt(slide, Inches(1.2), Inches(5.2), Inches(4), Inches(0.4),
        "Prepared for SRAM Executive Leadership  \u2022  March 2026", size=13, color=GRAY)

    txt(slide, Inches(1.2), Inches(6.5), Inches(6), Inches(0.3),
        "CRAWL  \u2192  WALK  \u2192  RUN", size=16, color=ACCENT_BLUE, bold=True)


def slide_exec_summary(prs):
    """Executive summary with the three-horizon framing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
        "SRAM has a rare structural advantage for AI \u2014 but needs a disciplined path to capture it",
        size=22, color=BLACK, bold=True)

    # Accent underline
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), fill=ACCENT_BLUE, border=None, radius=0.5)

    # Left: The Opportunity
    add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5), fill=BG_LIGHT, border=LIGHT_GRAY)

    txt(slide, Inches(1.1), Inches(1.6), Inches(5.2), Inches(0.4),
        "THE OPPORTUNITY", size=12, color=ACCENT_BLUE, bold=True)

    bullet_list(slide, Inches(1.1), Inches(2.1), Inches(5.2), Inches(4.5), [
        "$1B+ revenue with AXS ecosystem generating telemetry data no competitor matches",
        "No ERP, no centralized analytics \u2014 70% of dealer questions follow automatable patterns",
        "Trustpilot rating at 1.6/5.0 driven by support friction, not product quality",
        "VP Digital Products confirms: \"60-70% ready\" for bounded AI deployment",
        "Year 1 expected net value: $10.2M on $2.7M spend (3.8x return)",
        "Phased approach de-risks investment: 90-day pilot proves value before scaling",
    ], size=13)

    # Right: Three horizons preview
    add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(5.5), fill=BG_LIGHT, border=LIGHT_GRAY)

    txt(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.4),
        "THE FRAMEWORK: CRAWL \u2192 WALK \u2192 RUN", size=12, color=ACCENT_BLUE, bold=True)

    # Crawl box
    add_shape(slide, Inches(7.3), Inches(2.2), Inches(5.0), Inches(1.45), fill=CARD_WHITE, border=CRAWL_COLOR)
    phase_badge(slide, Inches(7.5), Inches(2.35), "CRAWL", CRAWL_COLOR)
    txt(slide, Inches(9.1), Inches(2.35), Inches(3), Inches(0.3),
        "Early Wins  \u2022  0\u201312 months", size=11, color=GRAY, bold=True)
    txt(slide, Inches(7.5), Inches(2.85), Inches(4.6), Inches(0.7),
        "Build AI literacy and workforce productivity.\nLow-risk tools for everyday work.", size=12, color=BODY)

    # Walk box
    add_shape(slide, Inches(7.3), Inches(3.85), Inches(5.0), Inches(1.45), fill=CARD_WHITE, border=WALK_COLOR)
    phase_badge(slide, Inches(7.5), Inches(4.0), "WALK", WALK_COLOR)
    txt(slide, Inches(9.1), Inches(4.0), Inches(3), Inches(0.3),
        "Strategic Implementation  \u2022  1\u20133 years", size=11, color=GRAY, bold=True)
    txt(slide, Inches(7.5), Inches(4.5), Inches(4.6), Inches(0.7),
        "Embed AI into manufacturing, supply chain,\nand operational core processes.", size=12, color=BODY)

    # Run box
    add_shape(slide, Inches(7.3), Inches(5.5), Inches(5.0), Inches(1.25), fill=CARD_WHITE, border=RUN_COLOR)
    phase_badge(slide, Inches(7.5), Inches(5.65), "RUN", RUN_COLOR)
    txt(slide, Inches(9.1), Inches(5.65), Inches(3), Inches(0.3),
        "Competitive Advantage  \u2022  3\u20135 years", size=11, color=GRAY, bold=True)
    txt(slide, Inches(7.5), Inches(6.1), Inches(4.6), Inches(0.5),
        "AI-powered products and ecosystem that\nno competitor can replicate.", size=12, color=BODY)


def slide_crawl(prs):
    """CRAWL phase detail"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), fill=CRAWL_COLOR, border=None, radius=0)
    txt(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.35),
        "CRAWL  \u2022  EARLY WINS", size=14, color=WHITE, bold=True)
    txt(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.45),
        "Improve workforce productivity and build AI literacy across the organization",
        size=20, color=WHITE, bold=True)

    # Time horizon + Goal bar
    txt(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.3),
        "TIME HORIZON: 0\u201312 MONTHS", size=11, color=CRAWL_COLOR, bold=True)
    txt(slide, Inches(4.5), Inches(1.3), Inches(6), Inches(0.3),
        "GOAL: Low-risk adoption that builds confidence and delivers immediate ROI",
        size=11, color=GRAY, bold=True)

    # --- 4 Use Case Cards ---
    card_w = Inches(2.85)
    card_h = Inches(3.8)
    start_x = Inches(0.8)
    card_y = Inches(1.8)
    gap = Inches(0.25)

    cases = [
        ("Engineering\nProductivity", CRAWL_COLOR,
         ["Summarize test results and\nengineering reports",
          "CAD exploration via natural\nlanguage queries",
          "Code review assistance for\nfirmware teams"],
         "GitHub Copilot\nClaude Code"),

        ("Operations\nProductivity", CRAWL_COLOR,
         ["Meeting summaries and\naction item extraction",
          "Supplier communication\ndrafting and translation",
          "Contract clause analysis\nand comparison"],
         "Microsoft Copilot\nNotion AI"),

        ("Customer\nService", CRAWL_COLOR,
         ["Dealer troubleshooting\nassistant (AXS/Hammerhead)",
          "AI-drafted support responses\nwith human review",
          "Knowledge base search\nand retrieval"],
         "Zendesk AI\nAmazon Bedrock"),

        ("Marketing\nProductivity", CRAWL_COLOR,
         ["Product copy generation\nfor catalogs and web",
          "Campaign brief drafting\nand A/B variants",
          "Social media content\ncreation and scheduling"],
         "ChatGPT Enterprise\nAdobe Firefly"),
    ]

    for i, (title, color, items, tools) in enumerate(cases):
        x = start_x + i * (card_w + gap)
        add_shape(slide, x, card_y, card_w, card_h, fill=CARD_WHITE, border=LIGHT_GRAY)

        # Icon area with color accent
        add_shape(slide, x, card_y, card_w, Inches(0.9), fill=BG_LIGHT, border=None, radius=0.04)
        txt(slide, x + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4), Inches(0.7),
            title, size=15, color=BLACK, bold=True)

        # Use cases
        for j, item in enumerate(items):
            txt(slide, x + Inches(0.2), card_y + Inches(1.05) + j * Inches(0.65),
                card_w - Inches(0.4), Inches(0.6),
                f"\u25B8  {item}", size=11, color=BODY)

        # Tools footer
        add_shape(slide, x, card_y + card_h - Inches(0.75), card_w, Inches(0.75),
                  fill=BG_LIGHT, border=None, radius=0.04)
        txt(slide, x + Inches(0.2), card_y + card_h - Inches(0.65),
            card_w - Inches(0.4), Inches(0.55),
            tools, size=10, color=ACCENT_BLUE, bold=True)

    # Bottom: Key Dependencies
    add_shape(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.1), fill=BG_LIGHT, border=LIGHT_GRAY)
    txt(slide, Inches(1.1), Inches(5.95), Inches(3), Inches(0.3),
        "KEY DEPENDENCIES", size=11, color=ACCENT_BLUE, bold=True)
    bullet_list(slide, Inches(1.1), Inches(6.25), Inches(11), Inches(0.6), [
        "Minimal technical dependency \u2014 these are SaaS tools, not infrastructure builds",
        "Security governance and data handling policies for AI tool usage",
        "Internal AI usage guidelines and acceptable use policy",
        "Employee training program (2\u20134 hours per function, not weeks)",
    ], size=11, color=BODY)


def slide_walk(prs):
    """WALK phase detail"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), fill=WALK_COLOR, border=None, radius=0)
    txt(slide, Inches(0.8), Inches(0.15), Inches(5), Inches(0.35),
        "WALK  \u2022  STRATEGIC IMPLEMENTATION", size=14, color=WHITE, bold=True)
    txt(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.45),
        "Embed AI into core operations to improve efficiency and margins",
        size=20, color=WHITE, bold=True)

    txt(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.3),
        "TIME HORIZON: 1\u20133 YEARS", size=11, color=WALK_COLOR, bold=True)
    txt(slide, Inches(4.5), Inches(1.3), Inches(6), Inches(0.3),
        "GOAL: Operational AI embedded in manufacturing, supply chain, and forecasting",
        size=11, color=GRAY, bold=True)

    # --- 4 Use Case Cards ---
    card_w = Inches(2.85)
    card_h = Inches(3.8)
    start_x = Inches(0.8)
    card_y = Inches(1.8)
    gap = Inches(0.25)

    cases = [
        ("Manufacturing\nQuality Inspection",
         ["Computer vision detecting\nmachining defects in real time",
          "Automated pass/fail grading\non carbon fiber layups",
          "Defect pattern analysis\nacross production runs"],
         "Azure Vision\nAWS Rekognition"),

        ("Predictive\nMaintenance",
         ["Sensor-based monitoring of\nfactory machine health",
          "Failure prediction before\nunplanned downtime",
          "Maintenance scheduling\noptimized by usage data"],
         "AWS IoT + SageMaker\nPTC ThingWorx"),

        ("Supply Chain\nForecasting",
         ["Demand forecasting using\nhistorical + seasonal data",
          "Inventory optimization\nacross OEM and dealers",
          "Stockout and markdown\nloss reduction"],
         "Databricks\nSnowflake + ML"),

        ("Data\nInfrastructure",
         ["Unified data layer connecting\nShopify, support, telemetry",
          "Integration of Quarq and\nHammerhead acquisition data",
          "Foundation for all Phase 3\nAI-powered products"],
         "Snowflake\nAWS Glue / dbt"),
    ]

    for i, (title, items, tools) in enumerate(cases):
        x = start_x + i * (card_w + gap)
        add_shape(slide, x, card_y, card_w, card_h, fill=CARD_WHITE, border=LIGHT_GRAY)

        add_shape(slide, x, card_y, card_w, Inches(0.9), fill=BG_LIGHT, border=None, radius=0.04)
        txt(slide, x + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4), Inches(0.7),
            title, size=15, color=BLACK, bold=True)

        for j, item in enumerate(items):
            txt(slide, x + Inches(0.2), card_y + Inches(1.05) + j * Inches(0.65),
                card_w - Inches(0.4), Inches(0.6),
                f"\u25B8  {item}", size=11, color=BODY)

        add_shape(slide, x, card_y + card_h - Inches(0.75), card_w, Inches(0.75),
                  fill=BG_LIGHT, border=None, radius=0.04)
        txt(slide, x + Inches(0.2), card_y + card_h - Inches(0.65),
            card_w - Inches(0.4), Inches(0.55),
            tools, size=10, color=ACCENT_BLUE, bold=True)

    # Key Dependencies
    add_shape(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.1), fill=BG_LIGHT, border=LIGHT_GRAY)
    txt(slide, Inches(1.1), Inches(5.95), Inches(3), Inches(0.3),
        "KEY DEPENDENCIES", size=11, color=ACCENT_BLUE, bold=True)
    bullet_list(slide, Inches(1.1), Inches(6.25), Inches(11), Inches(0.6), [
        "Clean operational data \u2014 manufacturing, inventory, and supply chain data must be structured and accessible",
        "Centralized data platform \u2014 Snowflake/Databricks unifying acquisition-era siloed systems",
        "ML engineers and data engineers \u2014 2\u20133 hires or contractor engagements required",
        "Integration with ERP/MES systems \u2014 the absence of an ERP is SRAM's biggest infrastructure gap",
    ], size=11, color=BODY)


def slide_run(prs):
    """RUN phase detail"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), fill=RUN_COLOR, border=None, radius=0)
    txt(slide, Inches(0.8), Inches(0.15), Inches(5), Inches(0.35),
        "RUN  \u2022  COMPETITIVE ADVANTAGE", size=14, color=WHITE, bold=True)
    txt(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.45),
        "Build differentiated AI-powered products and ecosystem",
        size=20, color=WHITE, bold=True)

    txt(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.3),
        "TIME HORIZON: 3\u20135 YEARS", size=11, color=RUN_COLOR, bold=True)
    txt(slide, Inches(4.5), Inches(1.3), Inches(7), Inches(0.3),
        "GOAL: AI integrated into SRAM hardware and digital ecosystem \u2014 creating an unassailable moat",
        size=11, color=GRAY, bold=True)

    # --- 4 Use Case Cards ---
    card_w = Inches(2.85)
    card_h = Inches(3.8)
    start_x = Inches(0.8)
    card_y = Inches(1.8)
    gap = Inches(0.25)

    cases = [
        ("Smart Drivetrain\nOptimization",
         ["AI-adjusted shift timing\nbased on terrain and cadence",
          "Predictive chain/cassette\nwear and replacement alerts",
          "Personalized gear ratios\nfor rider style and fitness"],
         "Edge AI (on-device)\nAXS Firmware SDK"),

        ("AI Rider\nCoaching",
         ["Unified performance dashboard:\npower, gearing, suspension, HR",
          "Real-time ride optimization\nand post-ride analysis",
          "Training plan generation\nfrom multi-sensor data"],
         "Hammerhead Karoo\nAXS Intelligence Platform"),

        ("Predictive\nMaintenance (B2C)",
         ["Component health monitoring\nfor riders and fleet operators",
          "Proactive dealer alerts:\nparts approaching end-of-life",
          "Subscription revenue from\ncontinuous monitoring"],
         "AXS Cloud\nAWS IoT Greengrass"),

        ("Ecosystem\nData Monetization",
         ["B2B anonymized telemetry\nfor OEMs and dealers",
          "Dealer intelligence: regional\ndemand and inventory signals",
          "Data-informed product\ndevelopment and testing"],
         "Snowflake Data Sharing\nAXS API Platform"),
    ]

    for i, (title, items, tools) in enumerate(cases):
        x = start_x + i * (card_w + gap)
        add_shape(slide, x, card_y, card_w, card_h, fill=CARD_WHITE, border=LIGHT_GRAY)

        add_shape(slide, x, card_y, card_w, Inches(0.9), fill=BG_LIGHT, border=None, radius=0.04)
        txt(slide, x + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4), Inches(0.7),
            title, size=15, color=BLACK, bold=True)

        for j, item in enumerate(items):
            txt(slide, x + Inches(0.2), card_y + Inches(1.05) + j * Inches(0.65),
                card_w - Inches(0.4), Inches(0.6),
                f"\u25B8  {item}", size=11, color=BODY)

        add_shape(slide, x, card_y + card_h - Inches(0.75), card_w, Inches(0.75),
                  fill=BG_LIGHT, border=None, radius=0.04)
        txt(slide, x + Inches(0.2), card_y + card_h - Inches(0.65),
            card_w - Inches(0.4), Inches(0.55),
            tools, size=10, color=ACCENT_BLUE, bold=True)

    # Key Dependencies
    add_shape(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.1), fill=BG_LIGHT, border=LIGHT_GRAY)
    txt(slide, Inches(1.1), Inches(5.95), Inches(3), Inches(0.3),
        "KEY DEPENDENCIES", size=11, color=RUN_COLOR, bold=True)
    bullet_list(slide, Inches(1.1), Inches(6.25), Inches(11), Inches(0.6), [
        "Mature data platform from Walk phase \u2014 unified telemetry, CRM, and inventory data",
        "Edge computing capability for on-device AI (AXS components, Hammerhead Karoo)",
        "Product management and firmware engineering capacity for AI-embedded hardware",
        "Rider privacy framework and GDPR-compliant data governance for telemetry monetization",
    ], size=11, color=BODY)


def slide_roadmap_visual(prs):
    """Visual timeline showing all three phases with key milestones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
        "The journey from productivity tools to competitive moat is sequenced and measurable",
        size=22, color=BLACK, bold=True)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), fill=ACCENT_BLUE, border=None, radius=0.5)

    # Timeline backbone
    lane_y = Inches(1.7)
    lane_h = Inches(1.5)
    labels_x = Inches(0.8)
    bar_x = Inches(3.2)
    total_bar_w = Inches(9.3)

    phases = [
        ("CRAWL", "Early Wins", "0\u201312 months", CRAWL_COLOR, 0, 0.24),
        ("WALK", "Strategic\nImplementation", "1\u20133 years", WALK_COLOR, 0.20, 0.60),
        ("RUN", "Competitive\nAdvantage", "3\u20135 years", RUN_COLOR, 0.56, 1.0),
    ]

    milestones_data = [
        # (phase_idx, rel_x within phase, label)
        [("Q1", "AI tools\ndeployed"), ("Q2", "Support\npilot live"), ("Q4", "Org-wide\nadoption")],
        [("Y1.5", "Data platform\nlive"), ("Y2", "MFG vision\nAI deployed"), ("Y3", "Forecasting\nat scale")],
        [("Y3.5", "Smart shifting\nships"), ("Y4", "AXS Intelligence\nlaunches"), ("Y5", "Data\nmonetization")],
    ]

    for i, (label, sublabel, time, color, start_frac, end_frac) in enumerate(phases):
        y = lane_y + i * (lane_h + Inches(0.15))

        # Phase label
        phase_badge(slide, labels_x, y + Inches(0.2), label, color, w=Inches(1.1), h=Inches(0.35))
        txt(slide, labels_x, y + Inches(0.65), Inches(1.8), Inches(0.6),
            sublabel, size=11, color=GRAY, bold=True)
        txt(slide, labels_x, y + Inches(1.15), Inches(1.8), Inches(0.3),
            time, size=10, color=GRAY)

        # Bar
        bx = bar_x + total_bar_w * start_frac
        bw = total_bar_w * (end_frac - start_frac)
        add_shape(slide, bx, y + Inches(0.15), bw, Inches(0.55), fill=color, border=None, radius=0.08)

        # Milestones on bar
        milestones = milestones_data[i]
        for j, (mtime, mlabel) in enumerate(milestones):
            mx = bx + bw * ((j + 0.5) / len(milestones))
            # Time label on bar
            txt(slide, mx - Inches(0.4), y + Inches(0.2), Inches(0.8), Inches(0.2),
                mtime, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            # Milestone label below bar
            txt(slide, mx - Inches(0.6), y + Inches(0.75), Inches(1.2), Inches(0.6),
                mlabel, size=9, color=BODY, align=PP_ALIGN.CENTER)

    # Year markers at bottom
    year_y = lane_y + 3 * (lane_h + Inches(0.15)) + Inches(0.1)
    add_shape(slide, bar_x, year_y, total_bar_w, Pt(2), fill=LIGHT_GRAY, border=None, radius=0)
    years = ["Today", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]
    for i, yr in enumerate(years):
        x = bar_x + total_bar_w * (i / (len(years) - 1)) - Inches(0.3)
        txt(slide, x, year_y + Inches(0.08), Inches(0.6), Inches(0.25),
            yr, size=10, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

    # Bottom callout
    add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.85), fill=BG_LIGHT, border=ACCENT_BLUE)
    txt(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.6),
        "\u201CThe 90-day pilot produces data to justify or redirect every subsequent investment. "
        "Bounded problems first. Broader transformation after measured proof.\u201D",
        size=13, color=BLACK, italic=True, font=FONT_SERIF)


def slide_financials(prs):
    """Financial impact summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
        "The investment case is strong across all scenarios, with downside bounded by pilot scope",
        size=22, color=BLACK, bold=True)
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), fill=ACCENT_BLUE, border=None, radius=0.5)

    # Financial table
    table_x = Inches(0.8)
    table_y = Inches(1.5)
    col_widths = [Inches(4.2), Inches(2.3), Inches(2.3), Inches(2.3)]
    row_height = Inches(0.55)

    headers = ["Category", "Conservative", "Expected", "Upside"]
    rows = [
        ("Support assistant savings", "$0.5M", "$1.6M", "$4.3M"),
        ("Forecasting & inventory savings", "$0.7M", "$1.6M", "$4.5M"),
        ("Documentation & translation", "$0.1M", "$0.24M", "$0.9M"),
        ("Customer retention revenue", "$1.5M", "$3.0M", "$4.5M"),
        ("Package size increase revenue", "$2.0M", "$4.0M", "$8.0M"),
        ("Proposal win-rate revenue", "$1.3M", "$2.5M", "$3.8M"),
    ]
    totals_rows = [
        ("Gross Value", "$6.1M", "$12.9M", "$26.0M"),
        ("Total Investment", "($1.6M)", "($2.7M)", "($3.5M)"),
    ]
    final_row = ("Net Year-1 Value", "$4.5M", "$10.2M", "$22.5M")
    roi_row = ("Return on Spend", "2.8x", "3.8x", "6.4x")

    # Header row
    for j, header in enumerate(headers):
        x = table_x + sum(w for w in [Inches(0)] + [col_widths[k] for k in range(j)])
        x = table_x
        for k in range(j):
            x += col_widths[k]
        add_shape(slide, x, table_y, col_widths[j], row_height, fill=BG_DARK, border=None, radius=0)
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        txt(slide, x + Inches(0.15), table_y + Inches(0.1), col_widths[j] - Inches(0.3), Inches(0.35),
            header, size=12, color=WHITE, bold=True, align=align)

    # Data rows
    for i, row in enumerate(rows):
        y = table_y + row_height * (i + 1)
        bg_color = BG_LIGHT if i % 2 == 0 else CARD_WHITE
        for j, cell in enumerate(row):
            x = table_x
            for k in range(j):
                x += col_widths[k]
            add_shape(slide, x, y, col_widths[j], row_height, fill=bg_color, border=None, radius=0)
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            txt(slide, x + Inches(0.15), y + Inches(0.1), col_widths[j] - Inches(0.3), Inches(0.35),
                cell, size=12, color=BODY, align=align)

    # Divider
    div_y = table_y + row_height * (len(rows) + 1)
    add_shape(slide, table_x, div_y, sum(col_widths), Pt(2), fill=ACCENT_BLUE, border=None, radius=0)

    # Totals
    for i, row in enumerate(totals_rows):
        y = div_y + Inches(0.05) + row_height * i
        for j, cell in enumerate(row):
            x = table_x
            for k in range(j):
                x += col_widths[k]
            add_shape(slide, x, y, col_widths[j], row_height, fill=BG_LIGHT, border=None, radius=0)
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            txt(slide, x + Inches(0.15), y + Inches(0.1), col_widths[j] - Inches(0.3), Inches(0.35),
                cell, size=12, color=BLACK, bold=True, align=align)

    # Final highlight row
    for final in [final_row, roi_row]:
        i_offset = 0 if final == final_row else 1
        y = div_y + Inches(0.05) + row_height * (len(totals_rows) + i_offset)
        fill_color = ACCENT_BLUE if final == final_row else BG_DARK
        for j, cell in enumerate(final):
            x = table_x
            for k in range(j):
                x += col_widths[k]
            add_shape(slide, x, y, col_widths[j], row_height, fill=fill_color, border=None, radius=0)
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            txt(slide, x + Inches(0.15), y + Inches(0.1), col_widths[j] - Inches(0.3), Inches(0.35),
                cell, size=14, color=WHITE, bold=True, align=align)

    # Right side: Key insight callout
    callout_x = Inches(0.8)
    callout_y = div_y + Inches(0.05) + row_height * (len(totals_rows) + 2) + Inches(0.3)
    add_shape(slide, callout_x, callout_y, Inches(11.1), Inches(0.85), fill=BG_LIGHT, border=CRAWL_COLOR)
    txt(slide, callout_x + Inches(0.3), callout_y + Inches(0.1), Inches(10.5), Inches(0.6),
        "Expected-case support savings: 120 support staff \u00d7 $90K avg cost \u00d7 60% AI-addressable "
        "\u00d7 35% productivity gain \u00d7 70% adoption = $1.6M. All projections use transparent, auditable math.",
        size=11, color=BODY, italic=True)


def slide_the_ask(prs):
    """Closing slide: The Decision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    txt(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.5),
        "THE ASK", size=14, color=ACCENT_BLUE, bold=True)
    txt(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.6),
        "Three decisions to begin SRAM's AI transformation",
        size=28, color=WHITE, bold=True)

    # Three decision cards
    decisions = [
        ("1", "Launch the 90-day\nsupport pilot",
         "Scope to AXS + Hammerhead dealer support. "
         "One product owner, one integration engineer. "
         "Human approval on all outputs. "
         "Kill criteria defined upfront: two consecutive "
         "weeks of quality drops triggers rollback.",
         "CRAWL", CRAWL_COLOR),

        ("2", "Begin data infrastructure\nplanning",
         "Initiate requirements for a unified data layer "
         "connecting Shopify, support, and telemetry. "
         "Assign an executive sponsor with budget "
         "authority. This is the foundation for "
         "Walk and Run phases.",
         "WALK", WALK_COLOR),

        ("3", "Hold sequencing\ndiscipline",
         "Bounded problems first. Broader transformation "
         "after measured proof. The 90-day pilot produces "
         "data to justify every subsequent investment. "
         "One failed pilot taught SRAM that moving too "
         "fast costs months of organizational buy-in.",
         "ALL PHASES", RUN_COLOR),
    ]

    card_w = Inches(3.6)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    card_y = Inches(2.0)

    for i, (num, title, desc, phase_label, color) in enumerate(decisions):
        x = start_x + i * (card_w + gap)

        add_shape(slide, x, card_y, card_w, Inches(4.5), fill=RGBColor(0x24, 0x34, 0x55), border=color)

        # Number circle
        circle = add_shape(slide, x + Inches(0.2), card_y + Inches(0.25),
                           Inches(0.5), Inches(0.5), fill=color, border=None, radius=0.5,
                           shape_type=MSO_SHAPE.OVAL)
        circle.text_frame.word_wrap = False
        p = circle.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

        txt(slide, x + Inches(0.9), card_y + Inches(0.25), card_w - Inches(1.1), Inches(0.7),
            title, size=17, color=WHITE, bold=True)

        txt(slide, x + Inches(0.25), card_y + Inches(1.2), card_w - Inches(0.5), Inches(2.8),
            desc, size=12, color=GRAY)

        phase_badge(slide, x + Inches(0.25), card_y + Inches(4.0), phase_label, color, w=Inches(1.6))

    # Bottom line
    txt(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
        "Expected Year-1 return: 3.8x  \u2022  Downside bounded by 90-day pilot  \u2022  "
        "Upside scales with SRAM's unique data moat",
        size=14, color=ACCENT_BLUE, bold=True)


# === MAIN ===

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_exec_summary(prs)
    slide_crawl(prs)
    slide_walk(prs)
    slide_run(prs)
    slide_roadmap_visual(prs)
    slide_financials(prs)
    slide_the_ask(prs)

    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "deliverables")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "AI_Adoption_Strategy_Roadmap.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
